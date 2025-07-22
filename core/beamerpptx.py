"""
Beamer QT
Copyright (C) 2025  Jorge Guerrero - acroper@gmail.com

This program is free software: you can redistribute it and/or modify
it under the terms of the GNU General Public License as published by
the Free Software Foundation, either version 3 of the License, or
(at your option) any later version.

This program is distributed in the hope that it will be useful,
but WITHOUT ANY WARRANTY; without even the implied warranty of
MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
GNU General Public License for more details.

You should have received a copy of the GNU General Public License
along with this program.  If not, see <http://www.gnu.org/licenses/>.
"""



from lxml import etree
import xml.etree.ElementTree as ET
import argparse
import re
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor


class BeamerPPT:
    
    def __init__(self):
        None
    
    
    def Export(self, Document):
        
        # Has to go through the document components
        
        self.prs = Presentation()
        
        aspect_ratio = Document.FrontMatter.AspectRatio
        if aspect_ratio == '43':
            self.prs.slide_width = Emu(9144000)
            self.prs.slide_height = Emu(6858000)
        else:
            self.prs.slide_width = Emu(12192000)
            self.prs.slide_height = Emu(6858000)
            
        
        # Create Title Page
        """Creates the title slide (this remains the same)."""
        slide_layout = self.prs.slide_layouts[0] # Title Slide
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = Document.FrontMatter.Title
        st_text = Document.FrontMatter.Subtitle
        author_text = Document.FrontMatter.Author.replace(r'\and', ', ')
        subtitle.text = f"{st_text}\n{author_text}".strip()
        
        
        
    